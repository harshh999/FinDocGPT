"""
Enhanced File Processing Module for FinDocGPT
Handles multiple file formats with robust error handling and content extraction
"""

import pandas as pd
import numpy as np
import pdfplumber
from docx import Document
from pptx import Presentation
import openpyxl
import streamlit as st
from io import BytesIO
import tempfile
import os
from typing import Tuple, Dict, Optional, Any

class AdvancedFileProcessor:
    """Advanced file processor with support for multiple financial document formats"""
    
    def __init__(self):
        self.supported_formats = {
            'text/plain': {'name': 'Text File', 'icon': '📄'},
            'application/pdf': {'name': 'PDF Document', 'icon': '📕'},
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': {'name': 'Word Document', 'icon': '📘'},
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': {'name': 'PowerPoint Presentation', 'icon': '📊'},
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': {'name': 'Excel Spreadsheet', 'icon': '📗'},
            'application/vnd.ms-excel': {'name': 'Excel Spreadsheet (Legacy)', 'icon': '📗'},
            'text/csv': {'name': 'CSV File', 'icon': '📋'},
            'application/json': {'name': 'JSON File', 'icon': '🔧'},
            'application/xml': {'name': 'XML File', 'icon': '🔧'},
            'text/xml': {'name': 'XML File', 'icon': '🔧'}
        }
    
    def get_file_info(self, uploaded_file) -> Dict[str, Any]:
        """Get comprehensive file information"""
        file_info = {
            'name': uploaded_file.name,
            'type': uploaded_file.type,
            'size_bytes': len(uploaded_file.getvalue()),
            'size_mb': len(uploaded_file.getvalue()) / (1024 * 1024),
            'format_info': self.supported_formats.get(uploaded_file.type, {'name': 'Unknown', 'icon': '❓'})
        }
        return file_info
    
    def process_file(self, uploaded_file) -> Tuple[Optional[str], Optional[Dict[str, Any]]]:
        """Main file processing method with comprehensive format support"""
        try:
            file_info = self.get_file_info(uploaded_file)
            file_type = uploaded_file.type
            
            # Display processing info
            icon = file_info['format_info']['icon']
            name = file_info['format_info']['name']
            st.info(f"{icon} Processing {name}: {file_info['name']} ({file_info['size_mb']:.2f} MB)")
            
            # Route to appropriate processor
            processors = {
                "text/plain": self._process_text,
                "application/pdf": self._process_pdf,
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document": self._process_docx,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation": self._process_pptx,
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": self._process_excel,
                "application/vnd.ms-excel": self._process_excel,
                "text/csv": self._process_csv,
                "application/json": self._process_json,
                "application/xml": self._process_xml,
                "text/xml": self._process_xml
            }
            
            processor = processors.get(file_type)
            if processor:
                content, metadata = processor(uploaded_file)
                if content:
                    metadata.update(file_info)
                    return content, metadata
                else:
                    st.error("Failed to extract content from file")
                    return None, None
            else:
                st.error(f"Unsupported file format: {file_type}")
                return None, None
                
        except Exception as e:
            st.error(f"Error processing file: {str(e)}")
            return None, None
    
    def _process_text(self, uploaded_file) -> Tuple[Optional[str], Optional[Dict]]:
        """Process plain text files"""
        try:
            content = str(uploaded_file.read(), "utf-8")
            metadata = {
                "file_type": "text",
                "encoding": "utf-8",
                "lines": len(content.split('\n')),
                "characters": len(content),
                "words": len(content.split())
            }
            return content, metadata
        except UnicodeDecodeError:
            try:
                uploaded_file.seek(0)
                content = str(uploaded_file.read(), "latin-1")
                metadata = {
                    "file_type": "text",
                    "encoding": "latin-1",
                    "lines": len(content.split('\n')),
                    "characters": len(content),
                    "words": len(content.split())
                }
                return content, metadata
            except Exception as e:
                st.error(f"Error reading text file: {str(e)}")
                return None, None
    
    def _process_pdf(self, uploaded_file) -> Tuple[Optional[str], Optional[Dict]]:
        """Process PDF files with enhanced extraction"""
        try:
            content = ""
            tables_found = 0
            images_found = 0
            
            with pdfplumber.open(uploaded_file) as pdf:
                page_count = len(pdf.pages)
                
                for i, page in enumerate(pdf.pages):
                    # Extract text
                    page_text = page.extract_text()
                    if page_text:
                        content += f"\n--- Page {i+1} ---\n"
                        content += page_text + "\n"
                    
                    # Extract tables
                    tables = page.extract_tables()
                    if tables:
                        tables_found += len(tables)
                        for j, table in enumerate(tables):
                            content += f"\n--- Table {j+1} on Page {i+1} ---\n"
                            for row in table:
                                if row:
                                    content += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
                    
                    # Count images (approximate)
                    if hasattr(page, 'images'):
                        images_found += len(page.images)
            
            metadata = {
                "file_type": "pdf",
                "pages": page_count,
                "tables": tables_found,
                "images": images_found,
                "characters": len(content),
                "words": len(content.split())
            }
            
            return content, metadata
            
        except Exception as e:
            st.error(f"Error processing PDF: {str(e)}")
            return None, None
    
    def _process_docx(self, uploaded_file) -> Tuple[Optional[str], Optional[Dict]]:
        """Process Word documents with comprehensive extraction"""
        try:
            doc = Document(uploaded_file)
            content = ""
            
            # Extract paragraphs
            paragraph_count = 0
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content += paragraph.text + "\n"
                    paragraph_count += 1
            
            # Extract tables
            table_count = len(doc.tables)
            for i, table in enumerate(doc.tables):
                content += f"\n--- Table {i+1} ---\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        content += row_text + "\n"
            
            # Extract headers and footers (if accessible)
            try:
                for section in doc.sections:
                    if section.header:
                        header_text = ""
                        for paragraph in section.header.paragraphs:
                            header_text += paragraph.text
                        if header_text.strip():
                            content = f"HEADER: {header_text}\n" + content
                    
                    if section.footer:
                        footer_text = ""
                        for paragraph in section.footer.paragraphs:
                            footer_text += paragraph.text
                        if footer_text.strip():
                            content += f"\nFOOTER: {footer_text}"
            except:
                pass  # Headers/footers might not be accessible
            
            metadata = {
                "file_type": "docx",
                "paragraphs": paragraph_count,
                "tables": table_count,
                "sections": len(doc.sections),
                "characters": len(content),
                "words": len(content.split())
            }
            
            return content, metadata
            
        except Exception as e:
            st.error(f"Error processing Word document: {str(e)}")
            return None, None
    
    def _process_pptx(self, uploaded_file) -> Tuple[Optional[str], Optional[Dict]]:
        """Process PowerPoint presentations with slide-by-slide extraction"""
        try:
            prs = Presentation(uploaded_file)
            content = ""
            slide_count = len(prs.slides)
            total_shapes = 0
            tables_found = 0
            
            for i, slide in enumerate(prs.slides, 1):
                content += f"\n=== SLIDE {i} ===\n"
                slide_shapes = 0
                
                for shape in slide.shapes:
                    slide_shapes += 1
                    
                    # Extract text from shapes
                    if hasattr(shape, "text") and shape.text.strip():
                        content += f"[Text Box]: {shape.text}\n"
                    
                    # Extract text from tables
                    if shape.has_table:
                        tables_found += 1
                        table = shape.table
                        content += f"[Table {tables_found}]:\n"
                        for row in table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells])
                            if row_text.strip():
                                content += row_text + "\n"
                        content += "\n"
                
                total_shapes += slide_shapes
                content += f"[Slide {i} contains {slide_shapes} elements]\n"
            
            metadata = {
                "file_type": "pptx",
                "slides": slide_count,
                "total_shapes": total_shapes,
                "tables": tables_found,
                "characters": len(content),
                "words": len(content.split())
            }
            
            return content, metadata
            
        except Exception as e:
            st.error(f"Error processing PowerPoint: {str(e)}")
            return None, None
    
    def _process_excel(self, uploaded_file) -> Tuple[Optional[str], Optional[Dict]]:
        """Process Excel files with comprehensive sheet analysis"""
        try:
            # Read all sheets
            excel_data = pd.read_excel(uploaded_file, sheet_name=None)
            content = ""
            total_rows = 0
            total_cols = 0
            numeric_cols = 0
            
            for sheet_name, df in excel_data.items():
                content += f"\n=== SHEET: {sheet_name} ===\n"
                
                if not df.empty:
                    sheet_rows, sheet_cols = df.shape
                    total_rows += sheet_rows
                    total_cols += sheet_cols
                    
                    # Column information
                    content += f"Dimensions: {sheet_rows} rows × {sheet_cols} columns\n"
                    content += f"Columns: {', '.join(df.columns.astype(str))}\n\n"
                    
                    # Numeric column analysis
                    numeric_columns = df.select_dtypes(include=[np.number]).columns
                    numeric_cols += len(numeric_columns)
                    
                    if len(numeric_columns) > 0:
                        content += "NUMERIC ANALYSIS:\n"
                        for col in numeric_columns:
                            stats = df[col].describe()
                            content += f"{col}: Mean={stats['mean']:.2f}, Std={stats['std']:.2f}, Range=[{stats['min']:.2f}, {stats['max']:.2f}]\n"
                        content += "\n"
                    
                    # Sample data (first 20 rows)
                    content += "SAMPLE DATA:\n"
                    sample_size = min(20, len(df))
                    for idx, row in df.head(sample_size).iterrows():
                        row_data = []
                        for col, val in row.items():
                            if pd.notna(val):
                                row_data.append(f"{col}: {val}")
                        if row_data:
                            content += " | ".join(row_data) + "\n"
                    
                    if len(df) > sample_size:
                        content += f"... ({len(df) - sample_size} more rows)\n"
                else:
                    content += "Empty sheet\n"
                
                content += "\n"
            
            metadata = {
                "file_type": "excel",
                "sheets": len(excel_data),
                "total_rows": total_rows,
                "total_columns": total_cols,
                "numeric_columns": numeric_cols,
                "characters": len(content),
                "words": len(content.split())
            }
            
            return content, metadata
            
        except Exception as e:
            st.error(f"Error processing Excel file: {str(e)}")
            return None, None
    
    def _process_csv(self, uploaded_file) -> Tuple[Optional[str], Optional[Dict]]:
        """Process CSV files with statistical analysis"""
        try:
            df = pd.read_csv(uploaded_file)
            content = f"CSV FILE ANALYSIS\n{'='*50}\n\n"
            
            rows, cols = df.shape
            content += f"Dimensions: {rows} rows × {cols} columns\n"
            content += f"Columns: {', '.join(df.columns)}\n\n"
            
            # Data types
            content += "DATA TYPES:\n"
            for col, dtype in df.dtypes.items():
                content += f"{col}: {dtype}\n"
            content += "\n"
            
            # Numeric analysis
            numeric_cols = df.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                content += "STATISTICAL SUMMARY:\n"
                for col in numeric_cols:
                    stats = df[col].describe()
                    content += f"\n{col}:\n"
                    content += f"  Count: {stats['count']}\n"
                    content += f"  Mean: {stats['mean']:.4f}\n"
                    content += f"  Std: {stats['std']:.4f}\n"
                    content += f"  Min: {stats['min']:.4f}\n"
                    content += f"  Max: {stats['max']:.4f}\n"
                content += "\n"
            
            # Sample data
            content += "SAMPLE DATA (First 30 rows):\n"
            sample_size = min(30, len(df))
            for idx, row in df.head(sample_size).iterrows():
                row_data = []
                for col, val in row.items():
                    if pd.notna(val):
                        row_data.append(f"{col}: {val}")
                if row_data:
                    content += " | ".join(row_data) + "\n"
            
            if len(df) > sample_size:
                content += f"\n... ({len(df) - sample_size} more rows)\n"
            
            metadata = {
                "file_type": "csv",
                "rows": rows,
                "columns": cols,
                "numeric_columns": len(numeric_cols),
                "missing_values": df.isnull().sum().sum(),
                "characters": len(content),
                "words": len(content.split())
            }
            
            return content, metadata
            
        except Exception as e:
            st.error(f"Error processing CSV file: {str(e)}")
            return None, None
    
    def _process_json(self, uploaded_file) -> Tuple[Optional[str], Optional[Dict]]:
        """Process JSON files"""
        try:
            import json
            data = json.load(uploaded_file)
            content = f"JSON FILE ANALYSIS\n{'='*50}\n\n"
            content += json.dumps(data, indent=2, ensure_ascii=False)
            
            metadata = {
                "file_type": "json",
                "structure": type(data).__name__,
                "characters": len(content),
                "words": len(content.split())
            }
            
            return content, metadata
            
        except Exception as e:
            st.error(f"Error processing JSON file: {str(e)}")
            return None, None
    
    def _process_xml(self, uploaded_file) -> Tuple[Optional[str], Optional[Dict]]:
        """Process XML files"""
        try:
            import xml.etree.ElementTree as ET
            
            content_bytes = uploaded_file.read()
            content_str = content_bytes.decode('utf-8')
            
            # Parse XML
            root = ET.fromstring(content_str)
            
            content = f"XML FILE ANALYSIS\n{'='*50}\n\n"
            content += f"Root Element: {root.tag}\n\n"
            content += "XML Content:\n"
            content += content_str
            
            metadata = {
                "file_type": "xml",
                "root_element": root.tag,
                "characters": len(content),
                "words": len(content.split())
            }
            
            return content, metadata
            
        except Exception as e:
            st.error(f"Error processing XML file: {str(e)}")
            return None, None